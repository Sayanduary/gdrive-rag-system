import io

import pymupdf
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.services.groq import GroqService


SUPPORTED_IMAGE_TYPES = {
    "image/jpeg",
    "image/png",
    "image/jpg",
    "image/webp",
}


groq = GroqService()


# ==================================================
# PDF
# ==================================================

def extract_pdf_page_text(page) -> str:
    """
    Extract native text from a PDF page.
    """

    try:
        return page.get_text("text").strip()
    except Exception as error:
        print(
            f"Native PDF text extraction failed: "
            f"{type(error).__name__}: {error}"
        )
        return ""


def page_needs_ocr(
    page,
    min_text_chars: int = 30,
) -> bool:
    """
    Decide whether a PDF page should be sent
    to vision OCR.

    Pages with very little native text are
    considered scanned/image pages.
    """

    text = extract_pdf_page_text(page)

    return len(text) < min_text_chars


def render_pdf_page(
    page,
    scale: float = 2.0,
) -> bytes:
    """
    Render a PDF page into PNG bytes.

    2x gives good OCR quality without creating
    unnecessarily enormous images.
    """

    print(
        f"Rendering PDF page at {scale}x..."
    )

    pixmap = page.get_pixmap(
        matrix=pymupdf.Matrix(
            scale,
            scale,
        ),
        alpha=False,
    )

    image = Image.frombytes(
        "RGB",
        (
            pixmap.width,
            pixmap.height,
        ),
        pixmap.samples,
    )

    buffer = io.BytesIO()

    image.save(
        buffer,
        format="PNG",
        optimize=True,
    )

    image_bytes = buffer.getvalue()

    print(
        f"Rendered image: "
        f"{pixmap.width}x{pixmap.height}, "
        f"{len(image_bytes)} bytes"
    )

    return image_bytes


def ocr_pdf_page(
    page,
    page_number: int,
) -> str:
    """
    OCR one scanned PDF page using Groq Vision.
    """

    print("=" * 60)
    print(
        f"OCR START - PDF PAGE {page_number}"
    )
    print("=" * 60)

    try:

        image_bytes = render_pdf_page(
            page=page,
            scale=2.0,
        )

        print(
            f"Sending page {page_number} "
            f"to Groq Vision..."
        )

        text = groq.ocr_image(
            image_bytes=image_bytes,
            mime_type="image/png",
        )

        text = (
            text.strip()
            if isinstance(text, str)
            else ""
        )

        print(
            f"OCR SUCCESS - PAGE {page_number}"
        )

        print(
            f"OCR characters: {len(text)}"
        )

        print("=" * 60)

        return text

    except Exception as error:

        print("=" * 60)
        print(
            f"OCR FAILED - PDF PAGE "
            f"{page_number}"
        )
        print(
            f"Error type: "
            f"{type(error).__name__}"
        )
        print(
            f"Error: {error}"
        )
        print("=" * 60)

        # IMPORTANT:
        # Do not kill the entire PDF because
        # one OCR page failed.
        return ""


def extract_pdf_text(
    file_bytes: bytes,
) -> str:
    """
    Extract complete PDF text.

    Native PDF text is preferred.

    Scanned pages are automatically sent
    to Groq Vision OCR.
    """

    document = pymupdf.open(
        stream=file_bytes,
        filetype="pdf",
    )

    pages = []

    total_pages = len(document)

    print("=" * 70)
    print(
        f"PDF PAGES: {total_pages}"
    )
    print("=" * 70)

    try:

        for page_number, page in enumerate(
            document,
            start=1,
        ):

            print()
            print(
                f"Processing PDF page "
                f"{page_number}/{total_pages}"
            )

            text = extract_pdf_page_text(
                page
            )

            # ==========================================
            # NORMAL TEXT PAGE
            # ==========================================

            if not page_needs_ocr(page):

                print(
                    f"Page {page_number}: "
                    f"NATIVE TEXT "
                    f"({len(text)} chars)"
                )

                if text:

                    pages.append(
                        f"PAGE {page_number}\n"
                        f"{text}"
                    )

                continue

            # ==========================================
            # SCANNED PAGE
            # ==========================================

            print(
                f"Page {page_number}: "
                f"SCANNED / LOW TEXT "
                f"({len(text)} chars)"
            )

            ocr_text = ocr_pdf_page(
                page=page,
                page_number=page_number,
            )

            if ocr_text:

                pages.append(
                    f"PAGE {page_number}\n"
                    f"{ocr_text}"
                )

            elif text:

                # Preserve any native text if
                # OCR returned nothing.
                print(
                    f"Page {page_number}: "
                    f"OCR empty; preserving "
                    f"native text"
                )

                pages.append(
                    f"PAGE {page_number}\n"
                    f"{text}"
                )

            else:

                print(
                    f"Page {page_number}: "
                    f"NO TEXT EXTRACTED"
                )

    finally:

        document.close()

    final_text = (
        "\n\n".join(pages)
        .strip()
    )

    print()
    print("=" * 70)
    print(
        f"FINAL PDF TEXT: "
        f"{len(final_text)} characters"
    )
    print("=" * 70)

    return final_text


# ==================================================
# IMAGE OCR
# ==================================================

def extract_image_text(
    file_bytes: bytes,
    mime_type: str,
) -> str:
    """
    OCR a standalone image.
    """

    print("=" * 60)
    print("IMAGE OCR START")
    print("=" * 60)

    try:

        print(
            f"Image MIME type: {mime_type}"
        )

        print(
            f"Image size: "
            f"{len(file_bytes)} bytes"
        )

        text = groq.ocr_image(
            image_bytes=file_bytes,
            mime_type=mime_type,
        )

        text = (
            text.strip()
            if isinstance(text, str)
            else ""
        )

        print(
            f"IMAGE OCR SUCCESS - "
            f"{len(text)} characters"
        )

        return text

    except Exception as error:

        print(
            "=" * 60
        )

        print(
            "IMAGE OCR FAILED"
        )

        print(
            f"Error type: "
            f"{type(error).__name__}"
        )

        print(
            f"Error: {error}"
        )

        print(
            "=" * 60
        )

        return ""


# ==================================================
# PPTX
# ==================================================

def extract_pptx_shape_text(
    shape,
) -> list[str]:

    extracted = []

    # ==========================================
    # GROUP
    # ==========================================

    if shape.shape_type == (
        MSO_SHAPE_TYPE.GROUP
    ):

        for child in shape.shapes:

            extracted.extend(
                extract_pptx_shape_text(
                    child
                )
            )

        return extracted

    # ==========================================
    # TEXT
    # ==========================================

    if hasattr(
        shape,
        "text",
    ):

        text = (
            shape.text
            .strip()
        )

        if text:

            extracted.append(
                text
            )

    # ==========================================
    # IMAGE
    # ==========================================

    if shape.shape_type == (
        MSO_SHAPE_TYPE.PICTURE
    ):

        try:

            image_bytes = (
                shape.image.blob
            )

            mime_type = (
                getattr(
                    shape.image,
                    "content_type",
                    None,
                )
                or "image/png"
            )

            print(
                "Processing PPTX image "
                "with Groq Vision..."
            )

            image_text = groq.ocr_image(
                image_bytes=image_bytes,
                mime_type=mime_type,
            )

            if image_text:

                extracted.append(
                    image_text.strip()
                )

        except Exception as error:

            print(
                "PPTX image OCR error: "
                f"{type(error).__name__}: "
                f"{error}"
            )

    return extracted


def extract_pptx_text(
    file_bytes: bytes,
) -> str:

    presentation = Presentation(
        io.BytesIO(file_bytes)
    )

    slides = []

    for slide_number, slide in enumerate(
        presentation.slides,
        start=1,
    ):

        slide_parts = []

        for shape in slide.shapes:

            slide_parts.extend(
                extract_pptx_shape_text(
                    shape
                )
            )

        if slide_parts:

            slides.append(
                f"SLIDE {slide_number}\n"
                + "\n".join(
                    slide_parts
                )
            )

    return (
        "\n\n".join(slides)
        .strip()
    )


# ==================================================
# MAIN PARSER
# ==================================================

def parse_file(
    file_bytes: bytes,
    file_name: str,
    mime_type: str,
) -> str:

    print()
    print("=" * 70)
    print("FILE PARSER")
    print("=" * 70)

    print(
        f"File: {file_name}"
    )

    print(
        f"MIME: {mime_type}"
    )

    print(
        f"Bytes: {len(file_bytes)}"
    )

    # ==================================================
    # PDF
    # ==================================================

    if mime_type == "application/pdf":

        text = extract_pdf_text(
            file_bytes
        )

        print(
            f"Final PDF text characters: "
            f"{len(text)}"
        )

        return text.strip()

    # ==================================================
    # IMAGE
    # ==================================================

    if mime_type in SUPPORTED_IMAGE_TYPES:

        text = extract_image_text(
            file_bytes=file_bytes,
            mime_type=mime_type,
        )

        print(
            f"Vision OCR characters: "
            f"{len(text)}"
        )

        return text.strip()

    # ==================================================
    # PPTX
    # ==================================================

    if mime_type == (
        "application/vnd.openxmlformats-officedocument"
        ".presentationml.presentation"
    ):

        text = extract_pptx_text(
            file_bytes
        )

        print(
            f"PPTX text characters: "
            f"{len(text)}"
        )

        return text.strip()

    # ==================================================
    # UNSUPPORTED
    # ==================================================

    raise ValueError(
        f"Unsupported file type: "
        f"{mime_type}"
    )